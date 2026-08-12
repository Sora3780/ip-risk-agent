"""발표 대본을 텍스트 파일로 뽑는다.

발표 중에 PowerPoint 발표자 보기를 못 쓰는 상황(공용 노트북, 미러링)을 대비해
대본만 따로 둔다. 메모장에서 열어도 깨지지 않도록 BOM 을 붙여 저장한다.

사용:  python scripts/export_script.py
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))

from apply_notes import lookup, stamp  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
BUILDER = ROOT / "scripts" / "build_deck.py"
OUT = ROOT.parent / "IP DeteDog 발표대본.txt"

WIDTH = 78

# 마무리 장을 뺀 구성이면 끝맺는 말이 없다. 대본에는 남겨 둔다.
CLOSING = """마지막 장을 띄운 채로 말한다. 별도 슬라이드는 없다.

정리하면 이렇습니다.
검사의 시작점을 사람에서 파일 변경으로 옮겼고,
등급은 규칙이 정하고 설명은 근거 기반 모델이 만들도록 역할을 나눴습니다.
특허는 초록을 기준으로 후보를 좁히고, 인용은 원문과 다시 대조해 검증했습니다.

이상입니다. 감사합니다. 질문 받겠습니다."""

QA = """질문이 나올 만한 지점과 답할 방향이다. 외워서 말하지 말고 요지만 기억할 것.

기존 라이선스 스캐너와 무엇이 다른가
  실행 시점이 다릅니다. 기존 도구는 사람이 실행할 때 돕니다.
  저희는 파일 변경 자체가 검사의 시작점입니다.

정확도를 신뢰할 수 있는가
  소규모 고정 테스트셋 기준의 회귀 결과입니다.
  일반적인 정확도라고 말할 수 있는 규모는 아니고, 코드 변경 시 기존 동작이
  깨졌는지 확인하는 용도입니다. 표본 확대는 후속 과제로 두고 있습니다.

특허 판정을 믿을 수 있는가
  초록 기준의 기술 유사도이며 침해 판단이 아닙니다.
  모델이 든 인용이 원문에 실제로 있는지 문자열로 다시 확인하고,
  없으면 그 근거는 버립니다. 전문가 검토가 필요한지도 함께 표시합니다.

왜 배포하지 않았는가
  로컬 폴더 감시가 로컬 실행을 전제로 합니다.
  Drive 변경 알림 방식으로 바꾸려면 배포가 필요하고,
  그때 서버 인증을 먼저 적용할 계획입니다.

LLM 이 등급을 정하는 것 아닌가
  아닙니다. 등급은 SPDX 식별자 기반 규칙이 정합니다.
  같은 입력에는 항상 같은 등급이 나오고 판정 이유가 코드로 추적됩니다.
  모델은 설명과 의무사항, 후속 조치만 만듭니다.

비용이나 API 한도는 괜찮은가
  변경이 없는 주기에는 외부 호출이 발생하지 않습니다.
  특허 검토에는 최소 10분 재검토 간격과 잔여 한도 확인을 두었고,
  월 한도에 도달하면 검사를 중단하고 사유를 안내합니다."""


def deck_path() -> Path:
    """빌더가 저장하는 경로를 그대로 따라간다."""
    src = BUILDER.read_text(encoding="utf-8")
    for var in ("OUT", "FALLBACK_OUT"):
        m = re.search(rf'^{var} = ROOT\.parent / "(.+?)"', src, re.M)
        if m:
            p = ROOT.parent / m.group(1)
            if p.exists():
                return p
    raise SystemExit("발표 파일을 찾지 못했다. 먼저 build_deck.py 를 실행할 것")


def title_of(slide) -> str:
    """슬라이드에서 가장 큰 글자를 제목으로 본다.

    빌더 소스의 순서로 맞추면, PowerPoint 에서 장을 지웠을 때 제목이 밀린다.
    파일에 있는 것만 보고 정한다.
    """
    best = (0.0, 1e9, "")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        top = shape.top or 0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                body = run.text.strip()
                size = run.font.size.pt if run.font.size else 0
                if not body or size < 16:
                    continue
                if (size, -top) > (best[0], -best[1]):
                    best = (size, top, body)
    return best[2] or "(제목 없음)"


def split_time(note: str) -> tuple[str, str]:
    lines = note.split("\n")
    if lines and lines[0].strip().startswith("(약"):
        return lines[0].strip(), "\n".join(lines[1:]).strip("\n")
    return "", note.strip("\n")


def main() -> int:
    prs = Presentation(str(deck_path()))

    parts: list[tuple[int, str, str, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = title_of(slide)
        # 대본 원본은 notes_text.py 다. pptx 를 PowerPoint 가 열고 있어
        # 아직 반영하지 못했더라도 최신 대본으로 뽑는다.
        fresh = lookup(title)
        note = stamp(fresh) if fresh else (
            slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else "")
        when, body = split_time(note)
        parts.append((i, title, when, body))

    buf: list[str] = []
    buf.append("IP DeteDog 발표 대본")
    buf.append("아주대 AI 부트캠프 · PBL 1차 MVP · 5조")
    buf.append("")
    buf.append("─" * WIDTH)
    buf.append("차례")
    buf.append("")
    for i, name, when, _ in parts:
        buf.append(f"  {i:02d}  {name:<34s}{when}")
    buf.append("")
    buf.append("  괄호 안은 그 장에서 말하는 데 걸리는 예상 시간이다.")
    buf.append("  분당 345자 기준이며, 실제로 읽어 보고 조절할 것.")
    buf.append("─" * WIDTH)

    for i, name, when, body in parts:
        buf.append("")
        buf.append("")
        buf.append("═" * WIDTH)
        head = f"{i:02d}  {name}"
        buf.append(f"{head}{' ' * max(1, WIDTH - len(head) - len(when))}{when}")
        buf.append("═" * WIDTH)
        buf.append("")
        buf.append(body)

    joined = "\n".join(b for *_x, b in parts)
    if "질문 받겠습니다" not in joined:
        for name, block in (("마무리 멘트", CLOSING), ("예상 질문 대비", QA)):
            buf.append("")
            buf.append("")
            buf.append("═" * WIDTH)
            buf.append(f"※  {name}")
            buf.append("═" * WIDTH)
            buf.append("")
            buf.append(block)

    buf.append("")
    buf.append("")
    buf.append("─" * WIDTH)
    buf.append("끝")

    OUT.write_text("\n".join(buf), encoding="utf-8-sig")
    chars = sum(len(re.sub(r"\s", "", b.split("[예상 질문 대비]")[0])) for *_x, b in parts)
    print(f"저장: {OUT}")
    print(f"  {len(parts)}장 · 본문 {chars}자 · 예상 {chars / 345:.0f}분")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
