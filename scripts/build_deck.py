"""IP DeteDog 발표 자료 생성.

구현 흐름과 기술 구성을 중심으로 한 발표용 자료.
원본 발표 자료는 보존하고 ``_수정본`` 파일을 새로 생성.

사용: python scripts/build_deck.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT.parent / "IP DeteDog 발표자료_수정본.pptx"
FALLBACK_OUT = ROOT.parent / "IP DeteDog 발표자료_향후계획반영.pptx"

W, H = 13.333, 7.5
M = 0.82

BG = RGBColor(0xF7, 0xF8, 0xFA)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x18, 0x22, 0x2D)
INK_2 = RGBColor(0x4D, 0x5A, 0x68)
INK_3 = RGBColor(0x83, 0x8E, 0x9A)
LINE = RGBColor(0xDF, 0xE4, 0xE9)
BLUE = RGBColor(0x2F, 0x6F, 0xD6)
BLUE_PALE = RGBColor(0xE9, 0xF1, 0xFF)
NAVY = RGBColor(0x0D, 0x21, 0x36)
RED = RGBColor(0xD9, 0x3D, 0x4B)
RED_PALE = RGBColor(0xFC, 0xEA, 0xEC)
ORANGE = RGBColor(0xE8, 0x8A, 0x20)
ORANGE_PALE = RGBColor(0xFE, 0xF3, 0xE4)
GREEN = RGBColor(0x13, 0x9B, 0x69)
GREEN_PALE = RGBColor(0xE6, 0xF5, 0xEF)
PURPLE = RGBColor(0x7A, 0x5A, 0xC8)

FAMILY = "Pretendard"


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def blank(prs: Presentation, fill=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = fill
    return slide


def _ea(run, name: str = FAMILY) -> None:
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def text(
    slide,
    x,
    y,
    w,
    h,
    blocks,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line=1.25,
    margin=0,
):
    """blocks: (본문, 크기, 색, 굵게[, 앞 간격 pt[, 글꼴]]) 목록."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    for i, block in enumerate(blocks):
        body, size, color, bold = block[:4]
        space = block[4] if len(block) > 4 else 0
        family = block[5] if len(block) > 5 else FAMILY
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        para.space_before = Pt(space)
        run = para.add_run()
        run.text = body
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = family
        _ea(run, family)
    return box


def rect(slide, x, y, w, h, fill, *, edge=None, radius=0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if edge:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def line_shape(slide, x, y, w, h=0.02, color=LINE):
    return rect(slide, x, y, w, h, color)


def dot(slide, x, y, d, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def footer(slide, page: int, source: str = "") -> None:
    line_shape(slide, M, 7.13, W - 2 * M, 0.012, LINE)
    text(slide, M, 7.2, 10.2, 0.18, [(source, 8.3, INK_3, False)])
    text(slide, 11.8, 7.18, 0.7, 0.18, [(f"{page:02d}", 9, INK_3, True)], align=PP_ALIGN.RIGHT)


def heading(slide, page: int, title_: str, sub: str | None = None, *, source: str = "") -> None:
    text(slide, M, 0.52, 0.55, 0.28, [(f"{page:02d}", 11, BLUE, True)])
    text(slide, M + 0.62, 0.47, 11.3, 0.55, [(title_, 28, INK, True)])
    if sub:
        text(slide, M + 0.62, 1.08, 11.1, 0.35, [(sub, 13, INK_3, False)])
    footer(slide, page, source)


def notes(slide, script: str) -> None:
    slide.notes_slide.notes_text_frame.text = script


def code_text(slide, x, y, w, h, body: str, size=10.5, color=INK_2, align=PP_ALIGN.LEFT):
    return text(slide, x, y, w, h, [(body, size, color, False, 0, "Consolas")], align=align)


# ---------------------------------------------------------------------------
# Slides


def s01_cover(prs: Presentation) -> None:
    slide = blank(prs, NAVY)
    rect(slide, 0, 0, 0.12, H, BLUE)
    text(slide, M, 0.72, 7.8, 0.32,
         [("아주대 AI 부트캠프 · PBL 1차 MVP · 5조", 12.5, RGBColor(0x9F, 0xAD, 0xBA), False)])
    text(slide, M, 1.76, 11.4, 0.95, [("IP DeteDog", 55, PAPER, True)])
    text(slide, M, 2.94, 11.2, 0.72,
         [("작업공간 변경 기반 IP 리스크 상시 점검", 27, PAPER, True)])
    line_shape(slide, M, 4.05, 1.35, 0.04, BLUE)
    text(slide, M, 4.39, 10.9, 0.9,
         [("로컬 폴더 · Google Drive 변경 감지", 15, RGBColor(0xB8, 0xC4, 0xCE), False),
          ("오픈소스 라이선스 · 선행 특허 초록 대조", 15, RGBColor(0xB8, 0xC4, 0xCE), False, 8),
          ("등급 · 원문 근거 · 후속 조치 제공", 15, RGBColor(0xB8, 0xC4, 0xCE), False, 8)], line=1.25)
    text(slide, M, 6.46, 10.2, 0.3,
         [("Evidence-grounded workspace IP risk agent", 10.5, RGBColor(0x6F, 0x83, 0x95), False)])
    notes(slide, "(약 35초)\n\n안녕하세요. 5조 발표를 맡은 OOO입니다.\n\n저희가 만든 도구 이름은 IP DeteDog 입니다.\n하는 일은 한 문장으로 이렇습니다.\n작업공간의 파일이 바뀌면, 그 시점에 라이선스와 선행 특허 위험을 대신 확인해 줍니다.\n\n여기서 IP 는 지식재산을 말합니다.\n저희가 다루는 건 두 가지입니다. 오픈소스 라이선스, 그리고 이미 등록된 특허입니다.\n\n감시 대상도 두 가지입니다. 내 PC 의 로컬 폴더, 그리고 팀이 함께 쓰는 구글 드라이브입니다.\n\n먼저 이걸 왜 만들었는지부터 말씀드리겠습니다.")


def s02_background(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 2, "프로젝트 제작 배경")
    text(slide, M + 0.62, 1.68, 7.0, 0.35, [("반복되는 작업물 변경", 14, BLUE, True)])
    text(slide, M + 0.62, 2.12, 6.9, 1.18,
         [("개발 중에는 계속 바뀌는 파일,", 27, INK, True),
          ("검토 시점에는 이미 늦어진 IP 확인", 27, RED, True, 5)], line=1.16)

    text(slide, 8.35, 1.72, 3.65, 0.28, [("기존 확인 흐름", 12, INK_3, True)])
    stages = [("파일 변경", BLUE), ("수동 검색", BLUE), ("배포 전 발견", RED)]
    for i, (label, color) in enumerate(stages):
        yy = 2.34 + i * 0.75
        dot(slide, 8.37, yy, 0.24, color)
        text(slide, 8.92, yy - 0.06, 2.8, 0.3, [(label, 14.5, color if i == 2 else INK, i == 2)])
        if i < 2:
            line_shape(slide, 8.48, yy + 0.34, 0.02, 0.43, LINE)

    line_shape(slide, M + 0.62, 3.74, 11.0, 0.02, LINE)
    issues = [
        ("수동 시작", "KIPRIS 검색어 구성과 패키지별 원문 확인", BLUE),
        ("정보 분산", "외부 라이브러리 · 소스 코드 · 프로젝트 설명 · 특허 정보의 분리", ORANGE),
        ("발견 지연", "기능 개발 이후의 교체 비용과 전문가 검토 범위 확대", RED),
    ]
    for i, (label, body, color) in enumerate(issues):
        yy = 4.16 + i * 0.63
        text(slide, M + 0.62, yy, 1.25, 0.28, [(label, 12.5, color, True)])
        text(slide, M + 2.02, yy - 0.03, 9.4, 0.31, [(body, 13.5, INK_2, False)])

    rect(slide, M + 0.62, 6.18, 11.0, 0.54, BLUE_PALE)
    text(slide, M + 0.92, 6.32, 10.45, 0.27,
         [("핵심 전환  ·  수동 검사 실행 → 파일 변경 기반 자동 점검", 14, INK, True)])
    notes(slide, "(약 1분 15초)\n\n이 프로젝트는 사고를 겪어서 시작한 게 아닙니다. 확인할 방법을 찾다가 시작했습니다.\n\n개발하는 동안 작업물은 계속 바뀝니다.\n라이브러리가 추가되고, 참고한 코드가 들어오고, 기획 문서가 고쳐집니다.\n그런데 IP 확인은 그 변화를 따라가지 않습니다.\n\n오른쪽 흐름을 봐 주십시오.\n파일이 바뀌고, 사람이 검색을 하고, 배포 직전에 발견합니다.\n가운데 단계가 문제입니다. 아무도 시작하지 않으면 확인은 일어나지 않습니다.\n\n이유를 세 가지로 정리했습니다.\n\n첫째, 수동 시작입니다.\n특허를 확인하려면 특허청이 운영하는 특허정보 검색 서비스인 키프리스에 들어가\n검색어를 직접 만들어야 하고, 라이선스는 패키지마다 원문을 찾아 읽어야 합니다.\n\n둘째, 정보 분산입니다. 라이브러리와 소스, 문서, 특허 정보가 서로 다른 곳에 있습니다.\n\n셋째, 발견 지연입니다. 기능이 다 만들어진 뒤에 알면 교체 비용이 훨씬 커집니다.\n\n찾아보니 도구가 없지는 않았습니다. 라이선스 검사 도구도 있고 특허 검색도 있습니다.\n다만 전부 사람이 실행할 때만 돕니다.\n그래서 저희는 검사의 시작점을 사람에서 파일 변경으로 옮겨 보기로 했습니다.")


def s03_target(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 3, "핵심 대상 및 사용 상황")
    rect(slide, M + 0.62, 1.62, 11.0, 0.75, NAVY)
    text(slide, M + 0.96, 1.82, 10.4, 0.32,
         [("핵심 대상  ·  전담 법무·IP 인력이 없는 초기 프로젝트 팀", 18, PAPER, True)])

    roles = [
        ("개발자", "의존성·소스 코드", "패키지 추가와 vendor 코드 유입 시\n라이선스 등급 및 의무사항 확인", BLUE),
        ("기획자 · PM", "기획·제안·명세 문서", "문서 변경 시 선행 특허 후보와\n겹치는 기술 구성 확인", PURPLE),
        ("팀 리더", "배포·검토 의사결정", "위험 위치·원문 근거·후속 조치의\n한 화면 확인", GREEN),
    ]
    for i, (role, scope, body, color) in enumerate(roles):
        x = M + 0.62 + i * 3.67
        text(slide, x, 2.78, 3.2, 0.34, [(role, 17, color, True)])
        text(slide, x, 3.27, 3.2, 0.3, [(scope, 12.5, INK, True)])
        line_shape(slide, x, 3.72, 3.25, 0.015, LINE)
        text(slide, x, 3.98, 3.25, 0.9, [(body, 12.5, INK_2, False)], line=1.38)

    line_shape(slide, M + 0.62, 5.02, 11.0, 0.02, LINE)
    flow = [("입력", "로컬 폴더 · 공유 Drive"), ("트리거", "파일 이벤트 · 주기 확인"),
            ("처리", "라이선스 · 특허 대조"), ("출력", "등급 · 근거 · 조치")]
    for i, (label, body) in enumerate(flow):
        x = M + 0.62 + i * 2.75
        text(slide, x, 5.38, 0.9, 0.25, [(label, 10.5, BLUE, True)])
        text(slide, x, 5.78, 2.35, 0.55, [(body, 13.5, INK, True)], line=1.25)
        if i < 3:
            text(slide, x + 2.32, 5.7, 0.28, 0.3, [("›", 20, INK_3, True)], align=PP_ALIGN.CENTER)
    notes(slide, "(약 45초)\n\n대상은 전담 법무나 IP 인력이 없는 초기 프로젝트 팀입니다.\n한 사람이 여러 역할을 겸하는 팀을 기준으로 잡았습니다.\n\n역할별로 확인해야 하는 것이 다릅니다.\n개발자는 패키지를 추가하거나 외부 코드를 가져올 때 라이선스 등급과 의무사항을 봐야 합니다.\n기획자와 PM 은 문서를 고칠 때 선행 특허와 겹치는 부분이 있는지 봐야 합니다.\n팀 리더는 배포를 결정할 때 위험 위치와 근거, 다음 조치를 한 번에 봐야 합니다.\n\n아래가 전체 흐름입니다.\n입력은 로컬 폴더와 공유 드라이브, 트리거는 파일 이벤트와 주기 확인,\n처리는 라이선스와 특허 대조, 출력은 등급과 근거와 조치입니다.\n\n이 흐름을 어떻게 구성했는지 다음 장에서 보겠습니다.")


def s04_architecture(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 4, "전체 동작 구조", "입력 출처의 차이를 changes 형식으로 통일한 단일 검사 파이프라인")
    cols = [
        ("01", "작업공간", "로컬 폴더\nGoogle Drive", BLUE_PALE, BLUE),
        ("02", "변경 감지", "로컬 파일 이벤트\nDrive 최종 수정 시각 비교", PAPER, BLUE),
        ("03", "입력 정규화", "파일 위치 · 변경 상태\n변경된 내용", PAPER, BLUE),
        ("04", "분석 엔진", "라이선스 규칙·RAG\nKIPRIS·Gemini", PAPER, PURPLE),
        ("05", "결과 제공", "사용자 화면의 위험 목록\n상세 근거·변경 이력", GREEN_PALE, GREEN),
    ]
    for i, (num, label, body, fill, color) in enumerate(cols):
        x = M + 0.25 + i * 2.42
        rect(slide, x, 1.92, 2.16, 2.18, fill, edge=LINE)
        text(slide, x + 0.2, 2.14, 0.45, 0.23, [(num, 10, color, True)])
        text(slide, x + 0.2, 2.55, 1.78, 0.34, [(label, 16, INK, True)])
        text(slide, x + 0.2, 3.12, 1.78, 0.62, [(body, 11.7, INK_2, False)], line=1.35)
        if i < 4:
            text(slide, x + 2.16, 2.83, 0.25, 0.3, [("›", 18, INK_3, True)], align=PP_ALIGN.CENTER)

    line_shape(slide, M + 0.25, 4.53, 11.84, 0.02, LINE)
    text(slide, M + 0.25, 4.85, 1.55, 0.28, [("외부 연동", 11.5, INK_3, True)])
    sources = [
        ("Google Drive API v3", "폴더 목록 · 수정시각 · 문서 export"),
        ("deps.dev · PyPI", "패키지 SPDX 식별 및 보완"),
        ("한국저작권위원회", "84종 라이선스 의무사항 원문"),
        ("KIPRIS Plus", "특허 후보 · 영문 초록 · 국문 명칭"),
        ("Google Gemini", "검색어 추출 · 근거 설명 · 초록 비교"),
    ]
    for i, (name, role) in enumerate(sources):
        yy = 5.24 + (i % 3) * 0.48
        xx = M + 0.25 + (0 if i < 3 else 6.15)
        text(slide, xx, yy, 2.05, 0.25, [(name, 11.3, BLUE if i != 3 else PURPLE, True)])
        text(slide, xx + 2.12, yy - 0.01, 3.6, 0.27, [(role, 11.3, INK_2, False)])
    notes(slide, "(약 1분 20초)\n\n전체 구조입니다. 다섯 단계로 되어 있습니다.\n\n여기서 중요한 건 세 번째, 입력 정규화 단계입니다.\n로컬 폴더와 구글 드라이브는 변경을 감지하는 방식이 완전히 다릅니다.\n로컬은 파일 이벤트를 받고, 드라이브는 수정 시각을 비교합니다.\n\n이 차이를 앞단에서 흡수했습니다.\n서버로 넘어갈 때는 파일 위치, 변경 유형, 변경 내용, 이 세 가지 형식으로 통일됩니다.\n그래서 검사 파이프라인은 하나만 있으면 됩니다.\n로컬이든 드라이브든 같은 함수로 들어갑니다.\n\n아래가 연동한 외부 서비스 다섯 개인데, 이름이 낯선 것들만 짚고 넘어가겠습니다.\n\n드라이브 API 는 구글이 제공하는 것으로, 파일 목록과 수정 시각을 가져오는 데 씁니다.\n\ndeps.dev 는 구글이 운영하는 오픈소스 패키지 정보 서비스입니다.\n패키지 이름과 버전을 주면 그 패키지의 라이선스를 알려줍니다.\n\n한국저작권위원회 데이터는 라이선스 84종의 의무사항 원문입니다. 공공데이터포털에서 받았습니다.\n\n키프리스는 앞서 말씀드린 특허 검색 서비스이고,\n제미나이는 구글의 생성형 AI 모델입니다. 저희는 검색어 생성과 설명 작성에 씁니다.\n\n이제 각 단계를 순서대로 보겠습니다. 먼저 감지 방식입니다.")


def s05_monitoring(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 5, "작업공간 감지 방식", "로컬 이벤트 감지와 Drive 주기 확인의 이원화")
    line_shape(slide, 6.66, 1.63, 0.02, 4.94, LINE)

    text(slide, M + 0.48, 1.72, 4.9, 0.35, [("로컬 폴더 · 실시간 이벤트", 18, BLUE, True)])
    local = [
        ("감지 도구", "watchdog 파일 감시 모듈 · 하위 폴더 전체 감시"),
        ("감지 이벤트", "파일 생성 · 수정 · 삭제 · 이동"),
        ("대상 필터", "의존성 설정 · 소스 코드 · 기획 문서 · 임시파일 제외"),
        ("호출 제어", "3초 디바운스 · 동일 저장 이벤트 일괄 처리"),
        ("전송 범위", "초기 전체 동기화 이후 변경 파일 본문만 전송"),
        ("전송 제한", "파일당 1MB 이하 · 변경 파일만 검사 서버 전달"),
    ]
    for i, (label, body) in enumerate(local):
        yy = 2.32 + i * 0.58
        text(slide, M + 0.48, yy, 1.28, 0.25, [(label, 11.2, INK_3, True)])
        text(slide, M + 1.83, yy - 0.02, 4.25, 0.3, [(body, 12.3, INK_2, i == 3)])

    text(slide, 7.15, 1.72, 4.9, 0.35, [("Google Drive · 기본 20초 폴링", 18, PURPLE, True)])
    drive = [
        ("인증", "서비스 계정 · 공유 폴더 뷰어 권한"),
        ("변경 비교", "파일 목록 API의 경로별 최종 수정 시각 비교"),
        ("무변경", "본문 미다운로드 · 저장된 findings 즉시 반환"),
        ("변경 발생", "일반 파일 download · Google 문서 text export"),
        ("폴더 범위", "공유 Drive 포함 · 최대 6단계 재귀 탐색"),
        ("서버 전달", "Drive 문서를 로컬 변경 파일과 동일한 입력 형식으로 변환"),
    ]
    for i, (label, body) in enumerate(drive):
        yy = 2.32 + i * 0.58
        text(slide, 7.15, yy, 1.28, 0.25, [(label, 11.2, INK_3, True)])
        text(slide, 8.5, yy - 0.02, 4.0, 0.3, [(body, 12.3, INK_2, i == 2)])

    rect(slide, M + 0.48, 6.04, 11.0, 0.56, NAVY)
    text(slide, M + 0.78, 6.19, 10.4, 0.26,
         [("공통 입력 정보  ·  파일 위치  |  변경 유형  |  변경 내용", 12.2, PAPER, True)],
         align=PP_ALIGN.CENTER)
    notes(slide, "(약 1분 25초)\n\n감지는 두 갈래로 나뉩니다.\n\n왼쪽이 로컬 폴더입니다.\nwatchdog 이라는 파이썬 라이브러리를 쓰는데, 폴더 안에서 파일이 생기거나 바뀌면\n그 사실을 바로 알려주는 도구입니다.\n\n여기서 문제가 하나 있었습니다.\n편집기가 한 번 저장할 때 이벤트를 여러 개 만듭니다. 임시 파일을 만들고, 본문을 쓰고,\n임시 파일을 지우는 식입니다. 그대로 두면 저장 한 번에 검사가 네다섯 번 돕니다.\n\n그래서 디바운스를 넣었습니다.\n마지막 이벤트로부터 3초 동안 조용하면 그때 한 번만 검사하는 방식입니다.\n저장이 계속되는 동안에는 기다리다가, 손을 떼면 그때 묶어서 처리합니다.\n\n전송은 변경된 파일 본문만, 파일당 1MB 이하로 제한했습니다.\n\n오른쪽이 구글 드라이브입니다. 여기는 파일 이벤트라는 게 없습니다.\n그래서 기본 20초마다 경로별 최종 수정 시각 목록만 비교합니다.\n변경이 없으면 본문을 아예 내려받지 않고 저장된 결과를 그대로 돌려줍니다.\n\n변경이 있을 때만 본문을 가져오는데, 구글 문서는 일반 파일이 아니라서 그냥 다운로드가 안 됩니다.\n그래서 텍스트로 변환해 받는 export 기능을 씁니다.\n공유 드라이브를 포함해 최대 6단계까지 하위 폴더를 봅니다.\n\n두 경로 모두 아래 세 가지 형식으로 정리되어 서버로 넘어갑니다.")


def s06_license(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 6, "라이선스 점검 절차", "비공개 상용 배포 기준의 규칙 판정과 문서 표기 대조")
    steps = [
        ("1", "의존성 확인", "프로젝트에 등록된\n외부 라이브러리 목록 추출", BLUE),
        ("2", "식별자 조회", "deps.dev 우선\nPyPI 정보로 누락 보완", BLUE),
        ("3", "소스 확인", "라이선스 표기 문구\n외부 코드의 출처 누락", ORANGE),
        ("4", "위험도 판정", "표준 식별자 규칙 적용\n가장 높은 위험 등급", RED),
        ("5", "표기 대조", "프로젝트 안내의 라이선스\n실제 사용 라이브러리와 비교", GREEN),
    ]
    for i, (num, label, body, color) in enumerate(steps):
        x = M + 0.25 + i * 2.42
        rect(slide, x, 1.78, 2.16, 2.25, PAPER, edge=LINE)
        rect(slide, x, 1.78, 2.16, 0.09, color)
        text(slide, x + 0.2, 2.08, 0.35, 0.23, [(num, 11, color, True)])
        text(slide, x + 0.2, 2.43, 1.75, 0.32, [(label, 15.5, INK, True)])
        text(slide, x + 0.2, 3.02, 1.78, 0.65, [(body, 11.4, INK_2, False)], line=1.38)
        if i < 4:
            text(slide, x + 2.16, 2.7, 0.25, 0.3, [("›", 18, INK_3, True)], align=PP_ALIGN.CENTER)

    text(slide, M + 0.25, 4.43, 2.5, 0.28, [("등급 기준", 12, INK_3, True)])
    tiers = [
        ("NOTICE", "MIT · BSD · Apache 등", GREEN_PALE, GREEN),
        ("REVIEW", "미식별 · 출처 미상", BG, INK_2),
        ("RESTRICTED", "LGPL · MPL · EPL 등", ORANGE_PALE, ORANGE),
        ("FORBIDDEN", "AGPL · GPL 등", RED_PALE, RED),
    ]
    for i, (tier, desc, fill, color) in enumerate(tiers):
        x = M + 0.25 + i * 2.82
        rect(slide, x, 4.86, 2.55, 0.83, fill, edge=LINE)
        text(slide, x + 0.2, 5.04, 2.1, 0.25, [(tier, 11.5, color, True)])
        text(slide, x + 0.2, 5.37, 2.1, 0.23, [(desc, 10.7, INK_2, False)])

    rect(slide, M + 0.25, 5.98, 11.83, 0.64, BLUE_PALE)
    text(slide, M + 0.55, 6.12, 11.25, 0.36,
         [("등급 판정  ·  표준 라이선스 식별자(SPDX) 기반 자체 규칙", 12.5, BLUE, True),
          ("근거 제공  ·  한국저작권위원회 84종 원문", 12.5, INK, True, 3)], line=1.08)
    notes(slide, "(약 1분 35초)\n\n라이선스 점검입니다. 판정 기준은 비공개 상용 배포입니다.\n회사가 소스를 공개하지 않고 제품을 파는 상황을 가정했다는 뜻입니다.\n\n한 갈래가 아니라 세 갈래로 봅니다.\n의존성에 등록된 외부 라이브러리, 소스 코드 안의 라이선스 표기, 그리고 문서에 적힌 라이선스입니다.\n\n2단계 식별자 조회를 보시면, 여기서 말하는 식별자는 SPDX 입니다.\n라이선스마다 붙어 있는 표준 이름으로, Apache-2.0 이나 AGPL-3.0 같은 형태입니다.\n사람이 쓰는 표현은 제각각이지만 이 식별자는 하나로 정해져 있어서 규칙을 만들 수 있습니다.\n\n여기에 보완을 하나 넣었습니다.\ndeps.dev 가 PyMuPDF 라는 PDF 라이브러리를 non-standard, 그러니까 표준 형식이 아니라고만 답합니다.\n그런데 실제로는 AGPL 입니다. 이 응답만 쓰면 최고 위험 등급이 확인 등급으로 내려갑니다.\n그래서 파이썬 패키지 저장소인 PyPI 에서 원문을 한 번 더 가져와 다시 해석합니다.\n\n등급 판정도 같은 이유로 바꿨습니다.\n처음에는 공공 API 가 주는 라이선스 판별 플래그를 쓰려고 했는데,\nGPL 2.0 을 위험으로 표시하지 못하고 ISC 는 과다 탐지했습니다.\n그래서 플래그 대신 SPDX 기반 자체 규칙을 만들었습니다. 조건이 겹치면 가장 높은 등급을 채택합니다.\n\n등급은 아래 네 단계이고, 한국어 원문은 등급 판정이 아니라 근거 인용에만 씁니다.")


def s07_rag(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 7, "라이선스 RAG 구성", "84종 원문 전체 입력이 아닌 관련 조항 Top 3 검색 방식",
            source="한국저작권위원회 오픈소스SW 라이선스정보 · gemini-embedding-001")
    stages = [
        ("근거 데이터", "라이선스 84종", "주요 특징\n배포 시 의무사항\n출처 정보", BLUE),
        ("청크 구성", "총 149개", "주요 특징 50\n의무사항 91\n개요 8", BLUE),
        ("임베딩", "768차원", "조각당 최대 320자\n벡터 크기 통일\n로컬 검색 인덱스", PURPLE),
        ("검색", "관련 근거 3건", "문제가 된 라이선스 우선\n의미 유사도 순위\n전체 원문 보완 검색", PURPLE),
        ("생성", "정해진 결과 항목", "위험 설명\n의무사항\n후속 조치", GREEN),
    ]
    for i, (label, metric, body, color) in enumerate(stages):
        x = M + 0.25 + i * 2.42
        text(slide, x, 1.78, 2.05, 0.28, [(label, 11.5, color, True)])
        text(slide, x, 2.2, 2.08, 0.48, [(metric, 23, INK, True)])
        line_shape(slide, x, 2.85, 1.92, 0.02, LINE)
        text(slide, x, 3.14, 1.98, 1.15, [(body, 12, INK_2, False)], line=1.43)
        if i < 4:
            text(slide, x + 2.05, 2.38, 0.3, 0.3, [("›", 19, INK_3, True)], align=PP_ALIGN.CENTER)

    rect(slide, M + 0.25, 4.68, 11.83, 0.64, NAVY)
    text(slide, M + 0.56, 4.86, 11.2, 0.26,
         [("규칙 판정 + 위험 위치 + 관련 근거 3건  →  Gemini  →  설명 · 의무사항 · 후속 조치", 11.5, PAPER, True)],
         align=PP_ALIGN.CENTER)

    text(slide, M + 0.25, 5.67, 1.55, 0.27, [("역할 분리", 11.5, INK_3, True)])
    text(slide, M + 1.82, 5.62, 4.2, 0.35, [("등급 결정", 13.5, BLUE, True), ("SPDX 규칙 엔진", 13.5, INK, False, 5)], line=1.08)
    text(slide, M + 5.08, 5.62, 4.2, 0.35, [("설명 생성", 13.5, PURPLE, True), ("RAG 근거 기반 Gemini", 13.5, INK, False, 5)], line=1.08)
    text(slide, M + 8.47, 5.62, 3.5, 0.35, [("근거 부족", 13.5, RED, True), ("근거 부족 표시 · 수동 검토 전환", 11.5, INK, False, 5)], line=1.08)
    notes(slide, "(약 1분 30초)\n\n근거를 어떻게 붙이는지입니다.\n\nRAG 는 검색 증강 생성이라고 하는데, 말이 어렵지 방식은 단순합니다.\n모델에게 그냥 물어보면 아는 대로 답하니까, 먼저 관련 문서를 찾아서 같이 넣어 주고\n그 문서를 근거로 답하게 하는 겁니다.\n\n라이선스 원문은 84종을 다 가지고 있습니다.\n그런데 이걸 통째로 넣으면 관련 없는 조항이 같이 들어가서 설명의 근거가 흐려집니다.\n\n그래서 149개 조항으로 나눴습니다. 주요 특징 50개, 의무사항 91개, 개요 8개입니다.\n조각당 최대 320자로 자르고 임베딩으로 색인했습니다.\n임베딩은 문장을 숫자 목록으로 바꾸는 것인데, 뜻이 비슷한 문장끼리 숫자도 가까워집니다.\n그래서 단어가 달라도 의미가 비슷한 조항을 찾을 수 있습니다.\n\n검색은 두 단계입니다.\n먼저 문제가 된 라이선스로 후보를 좁힙니다. AGPL 이 문제면 AGPL 조항만 봅니다.\n그다음 지금 상황을 문장으로 만들어 가장 가까운 조항 세 건을 고릅니다.\n\n그러면 아래처럼 됩니다.\n규칙이 정한 등급, 위험이 있는 위치, 관련 근거 세 건. 이 세 가지를 제미나이에 넘깁니다.\n모델은 설명과 의무사항과 후속 조치만 만듭니다.\n\n역할이 나뉘어 있다는 점을 강조하고 싶습니다.\n등급은 규칙이 정하고 설명은 모델이 만듭니다.\n같은 입력에 항상 같은 등급이 나오고, 왜 그 등급인지 코드로 추적됩니다.\n근거가 부족하면 그렇다고 표시하고 수동 검토로 넘깁니다.")


def s08_patent_search(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 8, "특허 후보 수집 방식", "문서 변화 확인부터 KIPRIS 초록 확보까지의 단계별 처리",
            source="KIPRIS Plus 영문 특허 정보 · 국문 특허 명칭 조회 서비스")
    rows = [
        ("01", "대상 문서 선별", "문서 제목·경로의 기획·제안·명세·설계 관련 표현", BLUE),
        ("02", "재검토 조건", "문서 내용 지문 변경 · 최소 10분 재검토 간격 · KIPRIS 잔여 한도", BLUE),
        ("03", "검색어 생성", "Gemini 핵심 아이디어 요약 · 영문 2~3단어 검색어 5~6개", PURPLE),
        ("04", "선행 특허 후보 검색", "KIPRIS 영문 특허 검색 · 검색어별 상위 5건 · 출원번호 중복 제거", PURPLE),
        ("05", "후보 우선순위", "여러 검색어에 반복 등장한 특허 우선 · 상위 6건 판정", ORANGE),
        ("06", "초록 상세 조회", "영문 초록 조회 API · 국문 특허 명칭 조회 API", GREEN),
    ]
    for i, (num, label, body, color) in enumerate(rows):
        yy = 1.66 + i * 0.79
        text(slide, M + 0.43, yy, 0.46, 0.25, [(num, 10.5, color, True)])
        text(slide, M + 1.05, yy - 0.03, 2.18, 0.3, [(label, 14, INK, True)])
        text(slide, M + 3.42, yy - 0.03, 8.2, 0.34, [(body, 12.7, INK_2, False)])
        if i < 5:
            line_shape(slide, M + 0.57, yy + 0.38, 0.02, 0.39, LINE)
        line_shape(slide, M + 1.05, yy + 0.53, 10.55, 0.012, LINE)

    rect(slide, M + 0.43, 6.48, 11.17, 0.34, ORANGE_PALE)
    text(slide, M + 0.68, 6.55, 10.68, 0.2,
         [("검색 후보 수 · 판정 수 · 미판정 수의 동시 기록", 11, ORANGE, True)])
    notes(slide, "(약 1분 20초)\n\n특허입니다. 라이선스와는 접근이 다릅니다.\n특허는 신규 등록이 계속 늘어나서 미리 색인해 둘 수가 없습니다. 그래서 실시간 검색으로 갑니다.\n\n여섯 단계인데 앞의 두 개가 비용을 막는 장치입니다.\n먼저 모든 문서를 보지 않습니다. 제목이나 경로에 기획, 제안, 명세 같은 표현이 있는 문서만 봅니다.\n그리고 내용이 실제로 바뀌었는지, 마지막 검토에서 10분이 지났는지,\n키프리스 사용 한도가 남았는지 확인합니다. 세 조건을 다 통과해야 검색이 시작됩니다.\n\n검색어는 제미나이가 만듭니다.\n키프리스 검색은 넣은 단어를 모두 포함하는 방식이라, 길게 넣으면 결과가 0건이 나옵니다.\n저희도 처음에 다섯 단어짜리 검색어를 넣었다가 아무것도 못 찾았습니다.\n그래서 영문 두세 단어짜리 짧은 검색어를 대여섯 개 만들도록 했습니다.\n\n검색 결과는 출원번호로 중복을 제거합니다.\n출원번호는 특허마다 붙는 고유 번호라서, 검색어가 달라도 같은 특허인지 알 수 있습니다.\n그다음 여러 검색어에 반복 등장한 특허를 앞에 두고, 상위 6건만 초록을 가져옵니다.\n초록은 특허 내용을 요약한 글입니다. 특허 전문은 너무 길어서 초록만 받습니다.\n\n전부 보지 않는다는 점은 숨기지 않습니다. 후보 수, 판정 수, 미판정 수를 같이 기록합니다.")


def s09_patent_compare(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 9, "특허 대조 및 검증 방식", "초록 기반 유사도 판정과 양쪽 원문 인용 검증")
    rect(slide, M + 0.43, 1.65, 3.3, 1.22, BLUE_PALE, edge=LINE)
    text(slide, M + 0.7, 1.88, 2.75, 0.26, [("입력 1 · 기획서", 12, BLUE, True)])
    text(slide, M + 0.7, 2.29, 2.75, 0.3, [("핵심 아이디어 + 원문", 14, INK, True)])
    rect(slide, M + 4.35, 1.65, 3.3, 1.22, ORANGE_PALE, edge=LINE)
    text(slide, M + 4.62, 1.88, 2.75, 0.26, [("입력 2 · 선행 특허", 12, ORANGE, True)])
    text(slide, M + 4.62, 2.29, 2.75, 0.3, [("출원번호 + 영문 초록", 14, INK, True)])
    rect(slide, M + 8.27, 1.65, 3.3, 1.22, GREEN_PALE, edge=LINE)
    text(slide, M + 8.54, 1.88, 2.75, 0.26, [("출력 · 정해진 결과 항목", 12, GREEN, True)])
    text(slide, M + 8.54, 2.27, 2.75, 0.5, [("유사도 · 겹침 · 차이\n인용 쌍 · 전문가 검토", 12.8, INK, True)], line=1.3)
    text(slide, 3.92, 2.08, 0.36, 0.3, [("+", 18, INK_3, True)], align=PP_ALIGN.CENTER)
    text(slide, 7.84, 2.08, 0.36, 0.3, [("→", 18, INK_3, True)], align=PP_ALIGN.CENTER)

    line_shape(slide, M + 0.43, 3.26, 11.17, 0.02, LINE)
    text(slide, M + 0.43, 3.58, 2.3, 0.3, [("비교 결과 구성", 12, INK_3, True)])
    schema = [
        ("유사도 등급", "높음 · 중간 · 낮음 · 무관"),
        ("비교 요약", "겹치는 기술 구성 · 기획서 고유 구성"),
        ("원문 근거", "기획서 인용 · 특허 초록 인용 · 일치 이유"),
        ("검토 표시", "전문가 검토 필요 여부 · 근거 충족 여부"),
    ]
    for i, (key, body) in enumerate(schema):
        yy = 4.03 + i * 0.48
        text(slide, M + 0.43, yy, 2.4, 0.24, [(key, 11.2, PURPLE, True)])
        text(slide, M + 2.82, yy - 0.02, 4.05, 0.28, [(body, 11.8, INK_2, False)])

    line_shape(slide, 7.03, 3.53, 0.02, 2.16, LINE)
    text(slide, 7.43, 3.58, 3.8, 0.3, [("인용 검증", 12, INK_3, True)])
    checks = [
        "공백 정규화 후 기획서 원문 포함 여부 확인",
        "공백 정규화 후 특허 초록 포함 여부 확인",
        "양쪽 모두 미존재한 인용 쌍 제거",
        "높음·중간 결과만 위험 목록 반영",
        "후보 대비 미판정 건수 및 범위 표시",
    ]
    for i, body in enumerate(checks):
        yy = 4.01 + i * 0.39
        dot(slide, 7.45, yy + 0.04, 0.12, GREEN if i < 3 else ORANGE)
        text(slide, 7.75, yy, 4.2, 0.25, [(body, 11.4, INK_2, False)])

    rect(slide, M + 0.43, 6.05, 11.17, 0.65, RED_PALE)
    text(slide, M + 0.72, 6.18, 10.6, 0.35,
         [("판정 범위  ·  특허 초록 기준 기술 유사도 및 조사 필요 지점", 12.7, RED, True),
          ("제외 범위  ·  청구항 침해 여부와 법적 결론", 12.7, INK, True, 3)], line=1.08)
    notes(slide, "(약 1분 20초)\n\n대조 방식입니다.\n\n제미나이에 넣는 건 두 가지입니다.\n왼쪽은 기획서의 핵심 아이디어와 원문, 가운데는 특허 출원번호와 영문 초록입니다.\n\n결과는 자유 서술이 아니라 정해진 항목으로 받습니다.\n유사도 등급, 겹치는 구성과 다른 구성, 인용 쌍, 전문가 검토 필요 여부입니다.\n항목을 고정해 두면 매번 같은 형식으로 오기 때문에 화면에 바로 붙일 수 있습니다.\n\n여기서 문제가 하나 있습니다. 모델이 인용을 지어낼 수 있다는 겁니다.\n이 문장이 겹칩니다, 라고 했는데 원문에 그런 문장이 없으면\n근거가 아니라 오히려 사람을 잘못된 판단으로 끌고 갑니다.\n\n그래서 오른쪽 검증을 넣었습니다.\n모델이 든 기획서 인용이 기획서 원문에 실제로 있는지, 특허 인용이 초록에 실제로 있는지\n글자 그대로 다시 찾아봅니다. 띄어쓰기 차이는 무시하고 비교합니다.\n양쪽 다 없으면 그 인용 쌍은 버립니다.\n그리고 유사도가 높음과 중간인 것만 위험 목록에 올립니다.\n\n마지막 줄이 중요합니다.\n이 결과는 초록 기준의 기술 유사도입니다. 청구항 침해 여부나 법적 결론이 아닙니다.\n청구항은 특허의 권리 범위를 정한 부분인데, 저희는 그걸 판단하지 않습니다.\n전문가에게 가져갈 지점을 좁혀 주는 것까지가 저희 범위입니다.")


def s10_delivery(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 10, "외부 API 연동 및 결과 제공 방식", "외부 데이터의 역할과 사용자 화면에 제공되는 정보")
    text(slide, M + 0.43, 1.63, 5.55, 0.3, [("외부 API와 역할", 13, BLUE, True)])
    services = [
        ("Google Drive API v3", "파일 목록 · 수정 시각 · Google 문서 텍스트 변환"),
        ("deps.dev · PyPI", "외부 라이브러리의 표준 라이선스 식별자 확인"),
        ("저작권위원회 공공데이터", "라이선스 84종의 한국어 의무사항과 출처"),
        ("KIPRIS Plus", "선행 특허 후보 · 영문 초록 · 국문 명칭"),
        ("Google Gemini API", "특허 검색어 · 근거 설명 · 기술 유사도 비교"),
    ]
    rect(slide, M + 0.43, 2.03, 5.6, 0.46, NAVY)
    text(slide, M + 0.65, 2.15, 2.08, 0.2, [("서비스", 9.8, PAPER, True)])
    text(slide, M + 2.82, 2.15, 2.9, 0.2, [("활용 정보", 9.8, PAPER, True)])
    for i, (service, role) in enumerate(services):
        yy = 2.49 + i * 0.64
        rect(slide, M + 0.43, yy, 5.6, 0.64, PAPER if i % 2 == 0 else BG)
        text(slide, M + 0.65, yy + 0.17, 2.05, 0.28, [(service, 10.5, BLUE if i < 3 else PURPLE, True)])
        text(slide, M + 2.82, yy + 0.13, 2.9, 0.38, [(role, 10.5, INK_2, False)], line=1.2)
        line_shape(slide, M + 0.43, yy + 0.62, 5.6, 0.012, LINE)
    rect(slide, M + 0.43, 5.89, 5.6, 0.72, BLUE_PALE)
    text(slide, M + 0.68, 6.05, 5.1, 0.38,
         [("분석 서버(FastAPI)  ·  입력 통합 · 분석 조정 · 진행 상태 · 결과 이력 관리", 10.7, BLUE, True)],
         align=PP_ALIGN.CENTER)

    line_shape(slide, 6.66, 1.63, 0.02, 5.02, LINE)
    text(slide, 7.13, 1.63, 4.95, 0.3, [("사용자 화면 제공 정보", 13, GREEN, True)])
    outputs = [
        ("위험 목록", "등급 · 위치 · 라이선스 · 유형"),
        ("상세 근거", "판정 사유 · 원문 출처 · 의무사항 · 권장 조치"),
        ("특허 상세", "기획서·초록 나란히 보기 · 검증된 인용 강조"),
        ("변경 상태", "새 위험 · 해소 위험 · 변동 없음"),
        ("진행 상황", "규칙 · RAG/LLM · 특허 단계별 진행률"),
        ("검사 범위", "미검사 파일 · 미판정 후보 · 근거 부족 표시"),
    ]
    for i, (label, body) in enumerate(outputs):
        yy = 2.17 + i * 0.63
        text(slide, 7.13, yy, 1.33, 0.25, [(label, 11.7, BLUE if i < 3 else INK_3, True)])
        text(slide, 8.55, yy - 0.02, 3.75, 0.31, [(body, 11.6, INK_2, False)])
        if i < 5:
            line_shape(slide, 7.13, yy + 0.39, 5.05, 0.012, LINE)

    rect(slide, 7.13, 6.08, 5.05, 0.55, GREEN_PALE)
    text(slide, 7.34, 6.21, 4.62, 0.24,
         [("저장 정보  ·  최신 결과 · 검사 이력 · 특허 재검토 상태", 10.5, GREEN, True)],
         align=PP_ALIGN.CENTER)
    notes(slide, "(약 55초)\n\n지금까지가 분석이고, 이 장은 연동과 결과 제공입니다.\n\n왼쪽이 외부 API 다섯 가지와 각각에서 무엇을 쓰는지입니다.\n아래 FastAPI 는 파이썬으로 서버를 만드는 도구입니다.\n이 서버가 로컬과 드라이브에서 온 입력을 합치고, 분석 순서를 조정하고,\n진행 상태와 검사 이력을 관리합니다.\n\n오른쪽이 사용자 화면에 실제로 보이는 정보입니다.\n위험 목록에 등급과 위치와 라이선스와 유형이 나오고,\n항목을 열면 판정 사유와 원문 출처, 의무사항, 권장 조치가 나옵니다.\n특허는 기획서와 초록을 나란히 놓고, 검증을 통과한 인용만 색으로 강조합니다.\n\n아래 세 개도 봐 주십시오.\n새로 생긴 위험인지 해소된 위험인지, 지금 어느 단계까지 진행됐는지,\n그리고 무엇을 못 봤는지를 같이 표시합니다.\n결과가 0건일 때 안전한 건지 볼 게 없었던 건지 구분되어야 하기 때문입니다.")


def s11_results(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 11, "검증 결과 및 처리 성능", "고정 테스트 자료 기반의 회귀 테스트와 연속 실행 측정",
            source="저장소 내 재현 가능한 회귀 테스트 기준")
    text(slide, M + 0.43, 1.65, 5.4, 0.3, [("테스트 구성", 13, BLUE, True)])
    tests = [
        ("라이선스", "숨긴 위험 6건 · 오탐 방지 5건"),
        ("사용 환경", "개인·스타트업·부트캠프·기획자·빈 폴더 5종"),
        ("특허", "명확·부분·애매 3건 · 정상 기획서 1건"),
    ]
    for i, (label, body) in enumerate(tests):
        yy = 2.12 + i * 0.72
        text(slide, M + 0.43, yy, 1.2, 0.25, [(label, 11.5, INK_3, True)])
        text(slide, M + 1.72, yy - 0.03, 4.3, 0.31, [(body, 12.3, INK, True)])
        line_shape(slide, M + 0.43, yy + 0.4, 5.45, 0.012, LINE)

    text(slide, M + 0.43, 4.43, 5.4, 0.3, [("연속 실행 측정", 13, GREEN, True)])
    perf = [("최초 전체 검사", "9.8초"), ("변경 없음", "0.0초 · 외부 호출 0건"), ("파일 1개 변경", "3.5초")]
    for i, (label, value) in enumerate(perf):
        yy = 4.89 + i * 0.51
        text(slide, M + 0.43, yy, 2.25, 0.25, [(label, 11.8, INK_2, False)])
        text(slide, M + 3.05, yy - 0.03, 2.55, 0.3, [(value, 13, GREEN, True)], align=PP_ALIGN.RIGHT)

    line_shape(slide, 6.66, 1.63, 0.02, 4.95, LINE)
    text(slide, 7.13, 1.65, 4.95, 0.3, [("현재 결과", 13, PURPLE, True)])
    metrics = [
        ("라이선스 재현율", "6 / 6", "숨긴 위험 기준"),
        ("라이선스 정밀도", "6 / 6", "탐지 6건 · 오탐 0건"),
        ("근거 인용", "5 / 6", "미상 라이선스 1건 제외"),
        ("페르소나 회귀", "5 / 5", "폴더 형태 5종"),
        ("특허 사례", "3 / 3", "정상 기획서 후보 0건"),
    ]
    for i, (label, value, note) in enumerate(metrics):
        yy = 2.14 + i * 0.73
        text(slide, 7.13, yy, 1.82, 0.25, [(label, 11.5, INK_2, False)])
        text(slide, 9.04, yy - 0.05, 0.95, 0.32, [(value, 15, GREEN, True)], align=PP_ALIGN.RIGHT)
        text(slide, 10.2, yy, 2.0, 0.25, [(note, 10.8, INK_3, False)])
        line_shape(slide, 7.13, yy + 0.4, 5.05, 0.012, LINE)

    rect(slide, M + 0.43, 6.38, 11.75, 0.34, BLUE_PALE)
    text(slide, M + 0.68, 6.45, 11.25, 0.2,
         [("해석 범위  ·  현재 저장소의 소규모 고정 테스트셋 기준 회귀 결과", 10.8, BLUE, True)])
    notes(slide, "(약 1분 20초)\n\n검증 결과입니다.\n\n먼저 테스트 구성입니다.\n라이선스는 정답을 미리 숨겨 둔 위험 6건과, 걸리면 안 되는 항목 5건으로 만들었습니다.\n사용 환경은 개인 개발, 스타트업, 부트캠프, 기획자 전용 폴더, 빈 폴더 다섯 가지입니다.\n특허는 명확한 경우, 부분적인 경우, 애매한 경우 세 건과 정상 기획서 한 건입니다.\n\n오른쪽이 결과입니다.\n라이선스는 숨긴 6건을 모두 찾았고 오탐은 없었습니다.\n근거 인용은 6건 중 5건입니다. 나머지 한 건은 라이선스가 식별되지 않은 항목인데,\n검색 조건 없이 84종을 뒤지면 엉뚱한 조항이 붙기 때문에 일부러 근거 없음으로 남겼습니다.\n특허는 세 사례를 모두 의도한 등급으로 구분했고, 정상 기획서에서는 후보가 0건이었습니다.\n\n왼쪽 아래가 처리 시간입니다.\n최초 전체 검사가 9.8초, 변경이 없을 때는 0.0초에 외부 호출 0건,\n파일 하나만 바꿨을 때는 3.5초입니다.\n\n한 가지 분명히 말씀드리면, 이건 소규모 고정 테스트셋 기준의 회귀 결과입니다.\n회귀 테스트는 코드를 고쳤을 때 기존 동작이 깨졌는지 확인하는 것을 말합니다.\n일반적인 정확도라고 말할 수 있는 규모는 아니고, 표본 확대는 후속 과제로 두고 있습니다.")


def s12_status(prs: Presentation) -> None:
    slide = blank(prs)
    heading(slide, 12, "현재 구현 현황", "현재 구현 현황 및 향후 계획")
    line_shape(slide, 6.66, 1.63, 0.02, 4.98, LINE)
    text(slide, M + 0.48, 1.72, 4.9, 0.34, [("구현 완료", 18, GREEN, True)])
    done = [
        "로컬 폴더 실시간 이벤트 감지 및 변경 파일 전송",
        "Google Drive 수정시각 폴링 및 Google 문서 export",
        "의존성·소스 헤더·문서 표기의 라이선스 교차 점검",
        "84종 라이선스 근거 RAG 및 Gemini 설명 생성",
        "KIPRIS 후보 검색·초록 비교·인용 검증",
        "사용자 화면의 위험 목록·상세 근거·진행률·변경 이력",
        "의존성·LLM·KIPRIS·문서 해시 캐시 및 한도 관리",
    ]
    for i, body in enumerate(done):
        yy = 2.31 + i * 0.56
        dot(slide, M + 0.49, yy + 0.05, 0.14, GREEN)
        text(slide, M + 0.82, yy, 5.35, 0.3, [(body, 11.8, INK_2, False)])

    text(slide, 7.15, 1.72, 4.9, 0.34, [("향후 계획", 18, BLUE, True)])
    future = [
        ("Drive 트리거 전환", "현재: 기본 20초 주기 확인\n목표: 배포 후 변경 알림 기반 즉시 검토"),
        ("변경 파일 증분 검토", "현재: 변경 발생 시 전체 작업공간 재검토\n목표: 최초 검사 이후 변경 파일만 검토"),
        ("공유 워크스페이스", "여러 사용자의 동일 작업공간 연결\n감시 대상·위험 목록·검사 이력 공동 관리"),
        ("파일 단위 최소 권한", "현재: 서비스 계정의 공유 폴더 전체 읽기\n목표: 선택 파일별 권한 부여와 접근 범위 최소화"),
    ]
    for i, (label, body) in enumerate(future):
        yy = 2.24 + i * 0.94
        dot(slide, 7.15, yy + 0.05, 0.14, BLUE)
        text(slide, 7.48, yy - 0.02, 4.75, 0.3, [(label, 13.2, INK, True)])
        text(slide, 7.48, yy + 0.35, 4.75, 0.52, [(body, 10.9, INK_2, False)], line=1.3)

    rect(slide, M + 0.48, 6.29, 11.0, 0.43, BLUE_PALE)
    text(slide, M + 0.75, 6.39, 10.5, 0.22,
         [("고도화 방향  ·  트리거 기반 검토 → 증분 처리 → 공동 관리 → 최소 권한", 11.4, BLUE, True)])
    notes(slide, "(약 1분 30초)\n\n현재 구현 현황과 앞으로의 계획입니다.\n\n왼쪽은 지금 동작하는 기능입니다. 감지부터 결과 제공까지 흐름은 끊기지 않고 이어집니다.\n\n오른쪽이 남은 과제인데, 앞의 두 개는 저희가 스스로 짚은 한계입니다.\n\n첫째, 드라이브는 지금 20초마다 확인하는 방식입니다.\n구글이 변경 알림을 보내주는 기능이 있긴 한데, 알림을 받으려면 외부에서 접속 가능한\n주소가 있어야 해서 배포가 전제됩니다.\n감지 부분은 따로 떼어 놨기 때문에 알림 받는 부분만 바꿔 끼우면 됩니다.\n\n둘째, 지금은 변경이 생기면 작업공간 전체를 다시 검사합니다.\n판정과 설명 단계에는 결과를 저장해 두는 캐시가 있어서 새로 판정된 항목만 모델을 부르지만,\n파일을 읽고 규칙을 돌리는 부분은 전체를 돕니다.\n최초 검사 이후에는 바뀐 파일만 보도록 바꾸는 게 다음 단계입니다.\n\n셋째와 넷째는 확장입니다.\n여러 사람이 하나의 작업공간을 같이 관리하는 것,\n그리고 공유 폴더 전체가 아니라 선택한 파일 단위로만 권한을 주는 것입니다.\n지금은 서비스 계정에 폴더 전체 읽기 권한을 주고 있어서, 이 범위를 좁히려고 합니다.\n\n정리하겠습니다.\n검사의 시작점을 사람에서 파일 변경으로 옮겼고,\n등급은 규칙이 정하고 설명은 근거 기반 모델이 만들도록 역할을 나눴습니다.\n특허는 초록을 기준으로 후보를 좁히고, 인용은 원문과 다시 대조해 검증했습니다.\n\n이상입니다. 감사합니다. 질문 받겠습니다.")


def s13_close(prs: Presentation) -> None:
    slide = blank(prs, NAVY)
    rect(slide, 0, 0, 0.12, H, BLUE)
    text(slide, M, 0.66, 7.5, 0.3, [("핵심 정리", 13, RGBColor(0x8F, 0xA1, 0xB2), True)])
    points = [
        ("1", "파일 변경 기반의 검사 자동 시작"),
        ("2", "규칙 판정과 RAG·Gemini 설명의 역할 분리"),
        ("3", "KIPRIS 초록 기반 선행 특허 후보 압축"),
        ("4", "위험 등급·원문 근거·후속 조치의 통합 제공"),
    ]
    for i, (num, body) in enumerate(points):
        yy = 1.43 + i * 0.88
        text(slide, M, yy, 0.55, 0.35, [(num, 17, BLUE, True)])
        text(slide, M + 0.8, yy - 0.02, 10.7, 0.45, [(body, 20, PAPER, True)])
    line_shape(slide, M, 5.26, 1.35, 0.04, BLUE)
    text(slide, M, 5.61, 10.9, 0.42,
         [("IP DeteDog  ·  작업공간 변경 기반 IP 리스크 상시 점검", 17, RGBColor(0xB6, 0xC2, 0xCD), False)])
    text(slide, M, 6.65, 5.0, 0.35, [("Q&A", 19, PAPER, True)])
    text(slide, 9.1, 6.68, 3.35, 0.25,
         [("아주대 AI 부트캠프 · 5조", 10.5, RGBColor(0x7E, 0x91, 0xA3), False)], align=PP_ALIGN.RIGHT)
    notes(slide, "(약 25초)\n\n정리하겠습니다.\n\n첫째, 검사의 시작점을 사람에서 파일 변경으로 옮겼습니다.\n둘째, 등급은 규칙이 정하고 설명은 근거 기반 모델이 만들도록 역할을 나눴습니다.\n셋째, 특허는 초록 기준으로 후보를 상위 6건까지 좁혔습니다.\n넷째, 위험 등급과 원문 근거와 다음 조치를 한 화면에서 제공합니다.\n\n이상입니다. 감사합니다. 질문 받겠습니다.\n\n[예상 질문 대비]\n· 기존 스캐너와 차이 → 실행 시점입니다. 저희는 파일 변경이 시작점입니다.\n· 정확도 신뢰성 → 고정 테스트셋 회귀 결과이고, 표본 확대는 후속 과제입니다.\n· 특허 판정 근거 → 초록 기준 기술 유사도이며 침해 판단이 아닙니다. 인용은 원문 대조로 검증합니다.\n· 배포 여부 → 로컬 폴더 감시가 로컬 실행을 전제로 합니다. Drive 알림 전환 시 서버 인증과 함께 배포합니다.")


def main() -> int:
    prs = deck()
    for build in (
        s01_cover,
        s02_background,
        s03_target,
        s04_architecture,
        s05_monitoring,
        s06_license,
        s07_rag,
        s08_patent_search,
        s09_patent_compare,
        s10_delivery,
        s11_results,
        s12_status,
        s13_close,
    ):
        build(prs)
    out = OUT
    try:
        prs.save(out)
    except PermissionError:
        out = FALLBACK_OUT
        prs.save(out)
        print(f"기존 수정본 사용 중 · 새 파일 저장: {out}")
    else:
        print(f"저장: {out}")
    print(f"슬라이드: {len(prs.slides)}장 · 발표자 노트 포함")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
