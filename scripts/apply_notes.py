"""notes_text.py 의 대본을 발표 자료에 반영한다.

빌더 소스와 pptx 양쪽에 넣는다.
pptx 를 PowerPoint 가 열고 있으면 그 부분만 건너뛰고 빌더는 갱신한다.

사용:  python scripts/apply_notes.py
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

sys.path.insert(0, str(Path(__file__).resolve().parent))

from notes_text import NOTES  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
BUILDER = ROOT / "scripts" / "build_deck.py"
CPM = 345

# 빌더의 제목과 대본 키가 다른 경우
ALIAS = {"현재 구현 현황": "최종 정리"}


def stamp(body: str) -> str:
    """맨 앞에 예상 시간을 붙인다. 대본을 읽어 보고 조절할 기준값이다."""
    chars = len(re.sub(r"\s", "", body))
    secs = int(round(chars / CPM * 60 / 5.0) * 5)
    if secs < 60:
        when = f"(약 {secs}초)"
    else:
        m, s = divmod(secs, 60)
        when = f"(약 {m}분)" if s == 0 else f"(약 {m}분 {s}초)"
    return f"{when}\n\n{body}"


def lookup(title: str) -> str | None:
    title = ALIAS.get(title, title)
    if title in NOTES:
        return NOTES[title]
    for key, body in NOTES.items():          # 앞부분만 같아도 받아들인다
        if title.startswith(key) or key.startswith(title):
            return body
    return None


def patch_builder() -> int:
    src = BUILDER.read_text(encoding="utf-8")
    done = 0
    for m in list(re.finditer(r"^def (s\d+_\w+)\(.*?(?=\ndef |\Z)", src, re.S | re.M)):
        block = m.group(0)
        hm = re.search(r'heading\(slide, \d+, "(.+?)"', block)
        title = hm.group(1) if hm else ("IP DeteDog" if "s01" in m.group(1) else "")
        body = lookup(title) if title else None
        if body is None:
            print(f"  건너뜀 · {m.group(1)} (대본 없음)")
            continue
        nm = re.search(r'    notes\(slide, ".*?"\)\n', block, re.S)
        if not nm:
            continue
        esc = stamp(body).replace("\\", "\\\\").replace('"', '\\"').replace("\n", "\\n")
        src = src.replace(block, block.replace(nm.group(0), f'    notes(slide, "{esc}")\n'))
        done += 1
    BUILDER.write_text(src, encoding="utf-8")
    return done


def deck_path() -> Path | None:
    src = BUILDER.read_text(encoding="utf-8")
    for var in ("OUT", "FALLBACK_OUT"):
        m = re.search(rf'^{var} = ROOT\.parent / "(.+?)"', src, re.M)
        if m and (ROOT.parent / m.group(1)).exists():
            return ROOT.parent / m.group(1)
    return None


def patch_deck(path: Path) -> int:
    from pptx import Presentation

    from export_script import title_of

    prs = Presentation(str(path))
    done = 0
    for slide in prs.slides:
        body = lookup(title_of(slide))
        if body is None:
            continue
        slide.notes_slide.notes_text_frame.text = stamp(body)
        done += 1
    prs.save(str(path))
    return done


def main() -> int:
    n = patch_builder()
    print(f"build_deck.py · 대본 {n}장 반영")

    path = deck_path()
    if path is None:
        print("발표 파일 없음 · pptx 반영 건너뜀")
        return 0
    try:
        with open(path, "ab"):
            pass
    except PermissionError:
        print(f"{path.name} · PowerPoint 가 열고 있어 건너뜀 (닫고 다시 실행할 것)")
        return 0
    print(f"{path.name} · 대본 {patch_deck(path)}장 반영")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
